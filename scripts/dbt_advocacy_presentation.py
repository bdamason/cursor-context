"""
Script: dbt_advocacy_presentation.py
Purpose: Generate a PowerPoint presentation advocating for dbt adoption during stack migration
Author: ESO BI Team
Date: 2025-12-04
Dependencies: python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ESO Brand Colors from style guide
ESO_NAVY = RGBColor(0x0d, 0x39, 0x4e)           # #0d394e - Primary brand color
HEALTH_BLUE_DK = RGBColor(0x00, 0x6d, 0xa8)     # #006da8 - Secondary
HEALTH_BLUE_LT = RGBColor(0x5e, 0xa4, 0xde)     # #5ea4de - Secondary light
ESO_GRAY_DK = RGBColor(0x6c, 0x81, 0x93)        # #6c8193 - Gray dark
ESO_GRAY_LT = RGBColor(0xf4, 0xf7, 0xf9)        # #f4f7f9 - Gray light (backgrounds)
CHART_ORANGE = RGBColor(0xeb, 0x70, 0x1e)       # #eb701e - Chart accent
CHART_TEAL = RGBColor(0x11, 0xa7, 0x96)         # #11a796 - Chart accent
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

# Logo path
LOGO_PATH = r"C:\cursor_repo\infra\kb\ESO\style\ESO_Enterprise_BI_Logo_White05224875928389294.png"


def create_presentation():
    """Create the dbt advocacy presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide1 = add_title_slide(
        prs,
        title="Do It Right the First Time",
        subtitle="Why dbt Should Be Part of Our Migration",
        notes="Opening: We have a unique opportunity with this migration. Let's talk about how we can maximize it."
    )
    
    # Slide 2: The Reality of Our Current State
    slide2 = add_content_slide(
        prs,
        title="The Reality of Our Current State",
        bullets=[
            "Business logic scattered across stored procedures",
            "Tribal knowledge - only certain people understand certain procs",
            "No automated testing - issues found in production",
            "Lineage is unknown - \"what breaks if I change this?\"",
            "Version control gaps - hard to track what changed when",
            "Key message: We all agree the current design needs improvement"
        ],
        notes="Establish common ground. This isn't criticism - it's recognition that we have an opportunity to improve. Ask: Do we agree these are real challenges?"
    )
    
    # Slide 3: Two Migration Paths
    slide3 = add_two_column_slide(
        prs,
        title="Two Migration Paths",
        left_header="Path A: Migrate, Then Fix Later",
        left_bullets=[
            "Lift-and-shift existing stored procedures",
            "Technical debt moves with us",
            "Eventually need to refactor anyway",
            "= Doing the work TWICE"
        ],
        right_header="Path B: Migrate Right the First Time",
        right_bullets=[
            "Use migration as opportunity to restructure",
            "Adopt dbt as part of migration scope",
            "Build it correctly from day one",
            "= Do the work ONCE"
        ],
        notes="This is the core choice. Path A feels safer but costs more. Path B requires upfront investment but delivers better ROI. Which makes more sense?"
    )
    
    # Slide 4: What is dbt?
    slide4 = add_content_slide(
        prs,
        title="What is dbt?",
        bullets=[
            "SQL-first transformation layer (\"the T in ELT\")",
            "Transforms data already in the warehouse",
            "Not replacing stored procedures entirely - providing structure",
            "Works with any modern platform (Snowflake, Databricks, Synapse)",
            "Open source with massive community and enterprise support"
        ],
        notes="dbt is not revolutionary technology - it's SQL with structure. If you can write a SELECT statement, you can write a dbt model. The magic is in the framework around it."
    )
    
    # Slide 5: Addressing the Orchestration Concern
    slide5 = add_content_slide(
        prs,
        title="Addressing the Orchestration Concern",
        subtitle="\"What about the app and controller table pattern?\"",
        bullets=[
            "Honest answer: The controller table pattern gets replaced",
            "But replaced with something better - dbt manages dependencies automatically",
            "ref('model_name') creates the dependency graph",
            "State tracked in manifest.json and run_results.json",
            "Key insight: Controller tables exist because stored procedures don't know about each other",
            "dbt models DO know about each other via ref() - the problem goes away"
        ],
        notes="This is the honest slide. Yes, the pattern changes. But the reason the controller table exists is because stored procedures are isolated. dbt models have built-in awareness of each other. We're not losing functionality - we're getting it natively."
    )
    
    # Slide 6: Addressing the Control Concern
    slide6 = add_content_slide(
        prs,
        title="Addressing the Control Concern",
        subtitle="\"I need fine-grained control over execution\"",
        bullets=[
            "dbt provides granular control:",
            "    dbt run --select model_name (run specific models)",
            "    dbt run --select tag:daily (run by tags)",
            "    dbt run --select +model_name (run with upstream deps)",
            "Model-level configurations for materialization, partitioning",
            "Pre/post hooks for custom SQL when needed",
            "You can still call stored procedures from dbt when truly necessary"
        ],
        notes="You actually get MORE control with dbt, not less. The selector syntax is incredibly powerful. And for edge cases where you truly need procedural logic, dbt can call stored procedures."
    )
    
    # Slide 7: Addressing the Performance Concern
    slide7 = add_content_slide(
        prs,
        title="Addressing the Performance Concern",
        subtitle="\"Stored procedures are faster\"",
        bullets=[
            "dbt compiles to pure SQL - zero runtime overhead",
            "The compiled SQL IS the stored procedure equivalent",
            "Incremental models for efficient processing (only new/changed rows)",
            "Platform-specific optimizations built into dbt adapters",
            "Same SQL = Same performance",
            "Reality: If your SELECT is slow, it will be slow in both"
        ],
        notes="This is a myth. dbt doesn't add any runtime overhead - it generates SQL and executes it. The performance is identical to hand-written SQL. Incremental models can actually improve performance."
    )
    
    # Slide 8: Security - Decoupling Logic from Data
    slide8 = add_two_column_slide(
        prs,
        title="Security - Decoupling Logic from Data",
        left_header="Stored Procedures = Logic IN Database",
        left_bullets=[
            "Need database write permissions to modify logic",
            "Harder to audit - need database access to review",
            "Rollback = database restore",
            "AI tools would need database access = security risk"
        ],
        right_header="dbt = Logic in Code, Separate from Data",
        right_bullets=[
            "Logic version-controlled in git - full audit trail",
            "Code reviews before changes reach production",
            "Analysts contribute without database write access",
            "dbt runs with service account - least privilege",
            "AI-safe: Cursor sees SQL logic, never actual data",
            "Zero data exposure when using AI assistance"
        ],
        notes="Security teams love this separation. You can give analysts access to contribute to transformations without giving them database write access. And critically - AI tools like Cursor work on code, not data. No security concerns."
    )
    
    # Slide 9: Why Analysts Need to Own This
    slide9 = add_content_slide(
        prs,
        title="Why Analysts Need to Own This",
        subtitle="Business knowledge should inform data organization",
        bullets=[
            "Analysts understand the business - they should define the models",
            "Current: Business logic locked in procs only engineers can modify",
            "With dbt: Analysts write SQL models, engineers review and deploy",
            "Business definitions documented WHERE they're implemented",
            "Shift from \"ticket to engineering\" to \"PR with review\""
        ],
        notes="This is about leveraging the right expertise. Analysts know what 'active customer' means in our business. They should be able to encode that definition, with engineering providing guardrails and review."
    )
    
    # Slide 10: Building the Knowledge Layer for AI
    slide10 = add_content_slide(
        prs,
        title="Building the Knowledge Layer for AI",
        subtitle="Future-proofing our data platform",
        bullets=[
            "Auto-generated documentation and column-level lineage",
            "dbt model relationships generate the knowledge graph (RAG)",
            "AI assistants understand our data - it's in structured YAML and SQL",
            "Stored procedures = opaque to AI",
            "dbt models with descriptions = AI-readable context",
            "\"Ask Steve, he wrote that proc\" becomes \"Ask the AI, it knows\""
        ],
        notes="This is forward-looking. AI tools are only as good as the context they have. dbt creates that context automatically. Stored procedures in a database are a black box to AI. dbt models are self-documenting."
    )
    
    # Slide 11: Testing as Code
    slide11 = add_content_slide(
        prs,
        title="Testing as Code",
        subtitle="Catch issues before production",
        bullets=[
            "Schema tests: unique, not_null, accepted_values, relationships",
            "Custom data tests written in SQL",
            "Source freshness monitoring",
            "Tests run BEFORE deployment - catch issues in CI, not production",
            "Current state: We find data issues when reports break",
            "With dbt: We prevent data issues from reaching reports"
        ],
        notes="How do we currently know when data is wrong? Usually when someone complains. dbt lets us define expectations and validate them automatically. Every deployment is tested."
    )
    
    # Slide 12: Parallel Work = Faster Delivery
    slide12 = add_content_slide(
        prs,
        title="Parallel Work = Faster Delivery",
        subtitle="We're not adding scope - we're adding contributors",
        bullets=[
            "Key insight: dbt enables parallel contributions",
            "Analysts can start defining models NOW - don't wait for infrastructure",
            "Engineers focus on pipeline/infrastructure",
            "Analysts focus on business logic models",
            "Without dbt: Sequential - engineers migrate, THEN someone fixes models",
            "With dbt: Parallel - infrastructure and modeling simultaneously",
            "More contributors = potentially SHORTER timeline"
        ],
        notes="This is the reframe. The concern is 'dbt adds scope and time.' The reality is dbt adds PEOPLE. Analysts can work in parallel. We're not extending the timeline - we're parallelizing it."
    )
    
    # Slide 13: Proposed Path Forward
    slide13 = add_content_slide(
        prs,
        title="Proposed Path Forward",
        bullets=[
            "Include dbt in migration project scope from day one",
            "Analysts begin modeling business logic in parallel with infrastructure",
            "Use dbt to define target state data models during migration",
            "Keep external orchestration for job scheduling",
            "Result: Clean, documented, tested, AI-ready platform on day one",
            "",
            "The question isn't 'Can we afford to include dbt?'",
            "It's 'Can we afford to migrate our technical debt and fix it twice?'"
        ],
        notes="Close with the call to action. The migration is happening regardless. The choice is whether we do it once or twice. I recommend we do it right the first time."
    )
    
    return prs


def add_title_slide(prs, title, subtitle, notes=""):
    """Add a title slide with centered title and subtitle."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    
    # Add background shape - ESO Navy
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = ESO_NAVY
    background.line.fill.background()
    
    # Add ESO logo in top left
    if os.path.exists(LOGO_PATH):
        logo = slide.shapes.add_picture(
            LOGO_PATH,
            Inches(0.5), Inches(0.4),
            height=Inches(0.8)
        )
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle - using Chart Orange for accent
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2),
        Inches(12.333), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = CHART_ORANGE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_content_slide(prs, title, bullets, subtitle="", notes=""):
    """Add a content slide with title and bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add header bar - ESO Navy
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.3)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ESO_NAVY
    header_bar.line.fill.background()
    
    # Add ESO logo in header
    if os.path.exists(LOGO_PATH):
        logo = slide.shapes.add_picture(
            LOGO_PATH,
            Inches(11.5), Inches(0.25),
            height=Inches(0.6)
        )
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(10.5), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Starting Y position for content
    content_y = 1.5
    
    # Add subtitle if provided - using Health Blue Dark
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(content_y),
            Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = HEALTH_BLUE_DK
        content_y = 2.2
    
    # Add bullet points
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(content_y),
        Inches(12.333), Inches(5.5 - (content_y - 1.5))
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        
        # Handle indented bullets (starting with spaces)
        if bullet.startswith("    "):
            para.text = bullet.strip()
            para.level = 1
            para.font.size = Pt(18)
        elif bullet.startswith("  "):
            para.text = bullet.strip()
            para.level = 1
            para.font.size = Pt(18)
        else:
            para.text = bullet
            para.level = 0
            para.font.size = Pt(20)
        
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(12)
        
        # Highlight key messages - using Chart Orange
        if bullet.startswith("Key") or "Key insight" in bullet or "Key message" in bullet:
            para.font.bold = True
            para.font.color.rgb = CHART_ORANGE
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def add_two_column_slide(prs, title, left_header, left_bullets, right_header, right_bullets, notes=""):
    """Add a two-column comparison slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add header bar - ESO Navy
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.3)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ESO_NAVY
    header_bar.line.fill.background()
    
    # Add ESO logo in header
    if os.path.exists(LOGO_PATH):
        logo = slide.shapes.add_picture(
            LOGO_PATH,
            Inches(11.5), Inches(0.25),
            height=Inches(0.6)
        )
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(10.5), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Left column header - ESO Gray Dark (represents "old way")
    left_header_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(5.8), Inches(0.6)
    )
    left_header_frame = left_header_box.text_frame
    left_header_para = left_header_frame.paragraphs[0]
    left_header_para.text = left_header
    left_header_para.font.size = Pt(22)
    left_header_para.font.bold = True
    left_header_para.font.color.rgb = ESO_GRAY_DK
    
    # Left column content
    left_content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2),
        Inches(5.8), Inches(4.8)
    )
    left_content_frame = left_content_box.text_frame
    left_content_frame.word_wrap = True
    
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            para = left_content_frame.paragraphs[0]
        else:
            para = left_content_frame.add_paragraph()
        para.text = bullet
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(10)
    
    # Right column header - Chart Teal (represents "better way")
    right_header_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.5),
        Inches(6), Inches(0.6)
    )
    right_header_frame = right_header_box.text_frame
    right_header_para = right_header_frame.paragraphs[0]
    right_header_para.text = right_header
    right_header_para.font.size = Pt(22)
    right_header_para.font.bold = True
    right_header_para.font.color.rgb = CHART_TEAL
    
    # Right column content
    right_content_box = slide.shapes.add_textbox(
        Inches(6.8), Inches(2.2),
        Inches(6), Inches(4.8)
    )
    right_content_frame = right_content_box.text_frame
    right_content_frame.word_wrap = True
    
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            para = right_content_frame.paragraphs[0]
        else:
            para = right_content_frame.add_paragraph()
        para.text = bullet
        para.font.size = Pt(18)
        para.font.color.rgb = DARK_TEXT
        para.space_after = Pt(10)
    
    # Add divider line - Health Blue Light
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.4), Inches(1.5),
        Inches(0.03), Inches(5.5)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = HEALTH_BLUE_LT
    divider.line.fill.background()
    
    # Add speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return slide


def main():
    """Main function to generate the presentation."""
    print("Creating dbt advocacy presentation...")
    
    prs = create_presentation()
    
    output_file = "dbt_advocacy_presentation.pptx"
    prs.save(output_file)
    
    print(f"Presentation saved to: {output_file}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

